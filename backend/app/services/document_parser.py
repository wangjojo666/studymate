from __future__ import annotations

from pathlib import Path


def parse_document(path: Path) -> list[tuple[int, str]]:
    # 按文件类型提取文本，统一返回 (页码, 文本) 列表供后续切片和检索使用。
    suffix = path.suffix.lower()
    if suffix == ".txt":
        return _parse_txt(path)
    if suffix == ".docx":
        return _parse_docx(path)
    if suffix == ".pptx":
        return _parse_pptx(path)
    if suffix == ".pdf":
        return _parse_pdf(path)
    raise ValueError("暂不支持该文件类型，请上传 PDF、PPTX、DOCX、TXT 或图片文件。")


def _parse_txt(path: Path) -> list[tuple[int, str]]:
    for encoding in ("utf-8", "utf-8-sig", "gbk"):
        try:
            return [(1, path.read_text(encoding=encoding))]
        except UnicodeDecodeError:
            continue
    return [(1, path.read_text(errors="ignore"))]


def _parse_docx(path: Path) -> list[tuple[int, str]]:
    try:
        from docx import Document as DocxDocument
    except ImportError as exc:
        raise RuntimeError("缺少 python-docx，无法解析 Word 文档。") from exc

    document = DocxDocument(str(path))
    paragraphs = [p.text.strip() for p in document.paragraphs if p.text.strip()]
    for table in document.tables:
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if cells:
                paragraphs.append(" | ".join(cells))
    return [(1, "\n".join(paragraphs))]


def _parse_pptx(path: Path) -> list[tuple[int, str]]:
    try:
        from pptx import Presentation
    except ImportError as exc:
        raise RuntimeError("缺少 python-pptx，无法解析 PPT 文档。") from exc

    presentation = Presentation(str(path))
    pages: list[tuple[int, str]] = []
    for index, slide in enumerate(presentation.slides, start=1):
        lines: list[str] = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                lines.append(shape.text.strip())
        pages.append((index, "\n".join(line for line in lines if line)))
    return pages


def _parse_pdf(path: Path) -> list[tuple[int, str]]:
    fitz_pages = _parse_pdf_with_pymupdf(path)
    if fitz_pages:
        return fitz_pages

    pypdf_pages = _parse_pdf_with_pypdf(path)
    if pypdf_pages:
        return pypdf_pages

    raise RuntimeError("当前 Python 环境缺少 PDF 文本解析库，请安装 pypdf 或 PyMuPDF 后重试。")


def _parse_pdf_with_pymupdf(path: Path) -> list[tuple[int, str]]:
    try:
        import fitz
    except ImportError:
        return []

    pages: list[tuple[int, str]] = []
    with fitz.open(str(path)) as document:
        for index, page in enumerate(document, start=1):
            pages.append((index, page.get_text("text")))
    return pages


def _parse_pdf_with_pypdf(path: Path) -> list[tuple[int, str]]:
    try:
        from pypdf import PdfReader
    except ImportError:
        try:
            from PyPDF2 import PdfReader
        except ImportError:
            return []

    reader = PdfReader(str(path))
    pages: list[tuple[int, str]] = []
    for index, page in enumerate(reader.pages, start=1):
        pages.append((index, page.extract_text() or ""))
    return pages
